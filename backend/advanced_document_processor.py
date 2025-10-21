import os
import io
import mimetypes
from typing import Dict, List, Tuple, Optional, Any
from dataclasses import dataclass
from datetime import datetime
import logging

# File processing libraries
import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup
from PIL import Image
import pytesseract

# Text analysis libraries
import langdetect
import textstat
import re

logger = logging.getLogger(__name__)

@dataclass
class DocumentMetadata:
    """Enhanced document metadata with analysis results"""
    filename: str
    file_type: str
    file_size: int
    page_count: int = 0
    word_count: int = 0
    character_count: int = 0
    language: str = "unknown"
    readability_score: float = 0.0
    content_quality: str = "unknown"
    has_images: bool = False
    has_tables: bool = False
    image_count: int = 0  # Added image_count
    table_count: int = 0  # Added table_count
    extracted_images: int = 0
    processing_time: float = 0.0
    content_categories: List[str] = None

@dataclass
class ProcessingResult:
    """Enhanced processing result with analysis"""
    text: str
    metadata: DocumentMetadata
    chunks: List[Dict[str, Any]] = None
    images: List[Dict[str, Any]] = None
    tables: List[Dict[str, Any]] = None
    errors: List[str] = None

class AdvancedDocumentProcessor:
    """Enhanced document processor with advanced text extraction and analysis"""
    
    SUPPORTED_TYPES = {
        'application/pdf': '.pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
        'text/plain': '.txt',
        'text/csv': '.csv',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
        'application/vnd.ms-excel': '.xls',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
        'text/html': '.html',
        'application/rtf': '.rtf',
        'image/png': '.png',
        'image/jpeg': '.jpg',
        'image/gif': '.gif',
        'image/bmp': '.bmp',
        'image/tiff': '.tiff'
    }
    
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    
    def __init__(self, enable_ocr: bool = True, extract_images: bool = True):
        self.enable_ocr = enable_ocr
        self.extract_images = extract_images
        self.setup_ocr()
    
    def setup_ocr(self):
        """Setup OCR configuration"""
        try:
            # Test if tesseract is available
            pytesseract.get_tesseract_version()
            self.ocr_available = True
            logger.info("OCR (Tesseract) is available")
        except Exception as e:
            self.ocr_available = False
            logger.warning(f"OCR not available: {e}")
    
    def get_supported_formats(self):
        """Return list of supported file formats"""
        return ['.pdf', '.docx', '.txt', '.csv', '.xlsx', '.xls', '.pptx', '.html', 
                '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']
        
    def is_format_supported(self, file_path_or_extension):
        """Check if file format is supported"""
        if '.' in file_path_or_extension:
            ext = os.path.splitext(file_path_or_extension)[1].lower()
        else:
            ext = file_path_or_extension.lower()
            if not ext.startswith('.'):
                ext = '.' + ext
        return ext in self.get_supported_formats()
    
    def validate_file(self, file_path: str, file_size: int) -> Tuple[bool, str]:
        """Enhanced file validation with detailed checks"""
        # Check file size
        if file_size > self.MAX_FILE_SIZE:
            return False, f"File size ({file_size:,} bytes) exceeds maximum limit ({self.MAX_FILE_SIZE:,} bytes)"
        
        # Check file extension
        file_ext = os.path.splitext(file_path)[1].lower()
        if not any(file_ext == ext for ext in self.SUPPORTED_TYPES.values()):
            return False, f"File type '{file_ext}' not supported. Supported types: {list(self.SUPPORTED_TYPES.values())}"
        
        # Check if file exists and is readable
        if not os.path.exists(file_path):
            return False, "File does not exist"
        
        try:
            with open(file_path, 'rb') as f:
                # Read first few bytes to verify file integrity
                header = f.read(1024)
                if len(header) == 0:
                    return False, "File appears to be empty"
        except Exception as e:
            return False, f"Cannot read file: {str(e)}"
        
        return True, "File is valid"
    
    def extract_text(self, file_path: str) -> ProcessingResult:
        """Enhanced text extraction with comprehensive analysis"""
        start_time = datetime.now()
        
        filename = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        file_ext = os.path.splitext(filename)[1].lower()
        
        # Initialize metadata
        metadata = DocumentMetadata(
            filename=filename,
            file_type=file_ext,
            file_size=file_size
        )
        
        errors = []
        text = ""
        images = []
        tables = []
        
        try:
            # Route to appropriate extraction method
            if file_ext == '.pdf':
                text, images, tables, meta_updates = self._extract_from_pdf(file_path)
            elif file_ext == '.docx':
                text, images, tables, meta_updates = self._extract_from_docx(file_path)
            elif file_ext == '.txt':
                text, meta_updates = self._extract_from_txt(file_path), {}
            elif file_ext in ['.csv']:
                text, tables, meta_updates = self._extract_from_csv(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                text, tables, meta_updates = self._extract_from_excel(file_path)
            elif file_ext == '.pptx':
                text, images, meta_updates = self._extract_from_pptx(file_path)
            elif file_ext == '.html':
                text, meta_updates = self._extract_from_html(file_path), {}
            elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']:
                text, meta_updates = self._extract_from_image(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")
            
            # Update metadata with extraction results
            for key, value in meta_updates.items():
                setattr(metadata, key, value)
            
            # Analyze extracted text
            self._analyze_text(text, metadata)
            
        except Exception as e:
            error_msg = f"Error extracting text from {filename}: {str(e)}"
            logger.error(error_msg)
            errors.append(error_msg)
        
        # Calculate processing time
        processing_time = (datetime.now() - start_time).total_seconds()
        metadata.processing_time = processing_time
        
        return ProcessingResult(
            text=text,
            metadata=metadata,
            images=images,
            tables=tables,
            errors=errors
        )
    
    def _extract_from_pdf(self, file_path: str) -> Tuple[str, List[Dict], List[Dict], Dict]:
        """Enhanced PDF extraction with image and table detection"""
        text_content = []
        images = []
        tables = []
        
        doc = fitz.open(file_path)
        page_count = len(doc)
        has_images = False
        
        for page_num in range(page_count):
            page = doc.load_page(page_num)
            
            # Extract text
            page_text = page.get_text()
            text_content.append(page_text)
            
            # Extract images if enabled
            if self.extract_images:
                image_list = page.get_images()
                if image_list:
                    has_images = True
                    for img_index, img in enumerate(image_list):
                        try:
                            # Get image data
                            xref = img[0]
                            pix = fitz.Pixmap(doc, xref)
                            
                            if pix.n - pix.alpha < 4:  # GRAY or RGB
                                img_data = pix.tobytes("png")
                                images.append({
                                    'page': page_num + 1,
                                    'index': img_index,
                                    'size': len(img_data),
                                    'width': pix.width,
                                    'height': pix.height
                                })
                            pix = None
                        except Exception as e:
                            logger.warning(f"Error extracting image {img_index} from page {page_num}: {e}")
            
            # Detect tables (simple heuristic)
            if self._detect_tables_in_text(page_text):
                tables.append({
                    'page': page_num + 1,
                    'content': self._extract_table_from_text(page_text)
                })
        
        doc.close()
        
        full_text = '\n'.join(text_content)
        
        return full_text, images, tables, {
            'page_count': page_count,
            'has_images': has_images,
            'has_tables': len(tables) > 0,
            'extracted_images': len(images)
        }
    
    def _extract_from_docx(self, file_path: str) -> Tuple[str, List[Dict], List[Dict], Dict]:
        """Enhanced DOCX extraction with table and image detection"""
        doc = docx.Document(file_path)
        text_content = []
        tables = []
        images = []
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            text_content.append(paragraph.text)
        
        # Extract tables
        for table_index, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            tables.append({
                'index': table_index,
                'rows': len(table_data),
                'columns': len(table_data[0]) if table_data else 0,
                'content': table_data
            })
        
        # Check for images (basic check)
        has_images = False
        try:
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    has_images = True
                    images.append({
                        'type': 'embedded',
                        'target': rel.target_ref
                    })
        except Exception as e:
            logger.warning(f"Error checking for images in DOCX: {e}")
        
        full_text = '\n'.join(text_content)
        
        return full_text, images, tables, {
            'has_images': has_images,
            'has_tables': len(tables) > 0,
            'extracted_images': len(images)
        }
    
    def _extract_from_txt(self, file_path: str) -> str:
        """Enhanced text file extraction with encoding detection"""
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        
        # If all encodings fail, read as binary and decode with errors='replace'
        with open(file_path, 'rb') as file:
            return file.read().decode('utf-8', errors='replace')
    
    def _extract_from_csv(self, file_path: str) -> Tuple[str, List[Dict], Dict]:
        """Extract text from CSV files"""
        try:
            # Try different encodings and separators
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                for sep in [',', ';', '\t', '|']:
                    try:
                        df = pd.read_csv(file_path, encoding=encoding, sep=sep)
                        break
                    except Exception:
                        continue
                else:
                    continue
                break
            else:
                raise ValueError("Could not parse CSV file")
            
            # Convert to text representation
            text_content = []
            headers = df.columns.tolist()
            text_content.append("Column Headers: " + ", ".join(headers))
            
            # Add sample rows (limit for performance)
            sample_size = min(100, len(df))
            for index, row in df.head(sample_size).iterrows():
                row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                text_content.append(row_text)
            
            if len(df) > sample_size:
                text_content.append(f"... and {len(df) - sample_size} more rows")
            
            # Create table representation
            table_data = {
                'rows': len(df),
                'columns': len(df.columns),
                'headers': headers,
                'sample_data': df.head(10).to_dict('records') if len(df) > 0 else []
            }
            
            return '\n'.join(text_content), [table_data], {'has_tables': True}
            
        except Exception as e:
            raise ValueError(f"Error processing CSV: {str(e)}")
    
    def _extract_from_excel(self, file_path: str) -> Tuple[str, List[Dict], Dict]:
        """Extract text from Excel files"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text_content = []
            tables = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                text_content.append(f"\n--- Sheet: {sheet_name} ---")
                
                if not df.empty:
                    headers = df.columns.tolist()
                    text_content.append("Headers: " + ", ".join([str(h) for h in headers]))
                    
                    # Add sample rows
                    sample_size = min(50, len(df))
                    for index, row in df.head(sample_size).iterrows():
                        row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                        text_content.append(row_text)
                    
                    # Store table metadata
                    tables.append({
                        'sheet': sheet_name,
                        'rows': len(df),
                        'columns': len(df.columns),
                        'headers': headers,
                        'sample_data': df.head(5).to_dict('records')
                    })
            
            return '\n'.join(text_content), tables, {'has_tables': True}
            
        except Exception as e:
            raise ValueError(f"Error processing Excel file: {str(e)}")
    
    def _extract_from_pptx(self, file_path: str) -> Tuple[str, List[Dict], Dict]:
        """Extract text from PowerPoint files"""
        try:
            prs = Presentation(file_path)
            text_content = []
            images = []
            slide_count = 0
            
            for slide_num, slide in enumerate(prs.slides):
                slide_count += 1
                slide_text = []
                slide_text.append(f"\n--- Slide {slide_num + 1} ---")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Check for images
                    if shape.shape_type == 13:  # Picture
                        images.append({
                            'slide': slide_num + 1,
                            'type': 'image'
                        })
                
                text_content.extend(slide_text)
            
            return '\n'.join(text_content), images, {
                'page_count': slide_count,
                'has_images': len(images) > 0,
                'extracted_images': len(images)
            }
            
        except Exception as e:
            raise ValueError(f"Error processing PowerPoint file: {str(e)}")
    
    def _extract_from_html(self, file_path: str) -> str:
        """Extract text from HTML files"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                html_content = file.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text and clean it up
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text
            
        except Exception as e:
            raise ValueError(f"Error processing HTML file: {str(e)}")
    
    def _extract_from_image(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from images using OCR"""
        if not self.enable_ocr or not self.ocr_available:
            return "", {}
        
        try:
            # Open and process image
            image = Image.open(file_path)
            
            # Perform OCR
            text = pytesseract.image_to_string(image)
            
            return text, {
                'has_images': True,
                'extracted_images': 1
            }
            
        except Exception as e:
            logger.warning(f"OCR failed for {file_path}: {e}")
            return "", {}
    
    def _detect_tables_in_text(self, text: str) -> bool:
        """Detect if text contains table-like structures"""
        lines = text.split('\n')
        
        # Look for common table indicators
        table_indicators = 0
        for line in lines:
            # Check for multiple spaces (column alignment)
            if len(re.findall(r'\s{3,}', line)) >= 2:
                table_indicators += 1
            # Check for tab characters
            if '\t' in line and len(line.split('\t')) >= 3:
                table_indicators += 1
            # Check for pipe separators
            if '|' in line and len(line.split('|')) >= 3:
                table_indicators += 1
        
        # If more than 20% of lines look like table rows
        return table_indicators > len(lines) * 0.2
    
    def _extract_table_from_text(self, text: str) -> List[List[str]]:
        """Extract table data from text (simple implementation)"""
        lines = text.split('\n')
        table_data = []
        
        for line in lines:
            if '|' in line:
                # Pipe-separated
                row = [cell.strip() for cell in line.split('|') if cell.strip()]
                if len(row) >= 2:
                    table_data.append(row)
            elif '\t' in line:
                # Tab-separated
                row = [cell.strip() for cell in line.split('\t') if cell.strip()]
                if len(row) >= 2:
                    table_data.append(row)
            elif len(re.findall(r'\s{3,}', line)) >= 2:
                # Space-separated (tricky, basic implementation)
                row = re.split(r'\s{3,}', line.strip())
                if len(row) >= 2:
                    table_data.append(row)
        
        return table_data
    
    def _analyze_text(self, text: str, metadata: DocumentMetadata):
        """Analyze text content for quality and characteristics"""
        if not text.strip():
            return
        
        # Basic statistics
        metadata.word_count = len(text.split())
        metadata.character_count = len(text)
        
        # Language detection
        try:
            metadata.language = langdetect.detect(text)
        except Exception:
            metadata.language = "unknown"
        
        # Readability analysis
        try:
            metadata.readability_score = textstat.flesch_reading_ease(text)
            
            # Content quality assessment
            if metadata.readability_score >= 90:
                metadata.content_quality = "very_easy"
            elif metadata.readability_score >= 80:
                metadata.content_quality = "easy"
            elif metadata.readability_score >= 70:
                metadata.content_quality = "fairly_easy"
            elif metadata.readability_score >= 60:
                metadata.content_quality = "standard"
            elif metadata.readability_score >= 50:
                metadata.content_quality = "fairly_difficult"
            elif metadata.readability_score >= 30:
                metadata.content_quality = "difficult"
            else:
                metadata.content_quality = "very_difficult"
                
        except Exception as e:
            logger.warning(f"Error calculating readability: {e}")
            metadata.readability_score = 0.0
            metadata.content_quality = "unknown"
        
        # Content categorization (basic keyword-based)
        metadata.content_categories = self._categorize_content(text)
    
    def _categorize_content(self, text: str) -> List[str]:
        """Basic content categorization based on keywords"""
        categories = []
        text_lower = text.lower()
        
        # Define category keywords
        category_keywords = {
            'technical': ['api', 'code', 'function', 'algorithm', 'database', 'software', 'programming'],
            'business': ['revenue', 'profit', 'market', 'customer', 'sales', 'strategy', 'business'],
            'legal': ['contract', 'agreement', 'legal', 'law', 'regulation', 'compliance', 'terms'],
            'academic': ['research', 'study', 'analysis', 'methodology', 'conclusion', 'abstract', 'references'],
            'financial': ['budget', 'cost', 'expense', 'investment', 'financial', 'accounting', 'money'],
            'medical': ['patient', 'diagnosis', 'treatment', 'medical', 'health', 'clinical', 'symptoms'],
            'educational': ['learn', 'teach', 'student', 'course', 'education', 'training', 'curriculum']
        }
        
        for category, keywords in category_keywords.items():
            keyword_count = sum(text_lower.count(keyword) for keyword in keywords)
            # If category keywords appear frequently, add to categories
            if keyword_count >= 3 or any(keyword in text_lower for keyword in keywords):
                categories.append(category)
        
        return categories if categories else ['general']
    
    def analyze_document(self, content: str, filename: str = "") -> dict:
        """Analyze document content and return comprehensive metadata"""
        # Create metadata object
        metadata = DocumentMetadata(
            filename=filename,
            file_type=os.path.splitext(filename)[1].lower() if filename else ".txt",
            file_size=len(content.encode('utf-8'))
        )
        
        # Perform analysis
        self._analyze_text(content, metadata)
        
        # Return as dictionary
        return {
            'word_count': metadata.word_count,
            'character_count': metadata.character_count,
            'language': metadata.language,
            'readability_score': metadata.readability_score,
            'content_quality': metadata.content_quality,
            'content_categories': metadata.content_categories,
            'has_images': metadata.has_images,
            'has_tables': metadata.has_tables,
            'image_count': metadata.image_count,
            'table_count': metadata.table_count
        }

# Create enhanced processor instance
def create_processor(enable_ocr: bool = True, extract_images: bool = True) -> AdvancedDocumentProcessor:
    """Factory function to create document processor"""
    return AdvancedDocumentProcessor(enable_ocr=enable_ocr, extract_images=extract_images)